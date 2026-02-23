"""PowerPoint reader tool using python-pptx."""

from __future__ import annotations

import os
from typing import Optional

from google.adk.tools import FunctionTool


def pptx_reader(file_path: str, slide_number: Optional[int] = None) -> str:
    """Read PowerPoint presentations using python-pptx.

    Inserts [Slide N] markers for each slide.
    """
    if not os.path.exists(file_path):
        return f"[ERROR] File not found: {file_path}"

    try:
        from pptx import Presentation
    except ImportError:
        return "[ERROR] python-pptx not installed. Run: pip install python-pptx"

    try:
        prs = Presentation(file_path)
        total_slides = len(prs.slides)

        output = [f"PowerPoint: {os.path.basename(file_path)}"]
        output.append(f"Total slides: {total_slides}\n")

        for slide_idx, slide in enumerate(prs.slides, 1):
            # Skip if specific slide requested and this isn't it
            if slide_number is not None and slide_idx != slide_number:
                continue

            output.append(f"[Slide {slide_idx}]")

            # Extract text from all shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                output.append("\n".join(slide_text))
            else:
                output.append("(No text content)")

            output.append("")  # Blank line between slides

        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Failed to read PowerPoint: {str(e)}"


pptx_reader_tool = FunctionTool(pptx_reader)
